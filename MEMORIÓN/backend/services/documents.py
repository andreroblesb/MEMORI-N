from __future__ import annotations

import json
import re
from abc import ABC, abstractmethod
from pathlib import Path

from pydantic import BaseModel, Field


class DocumentChunkRequest(BaseModel):
    path: str = Field(min_length=1, max_length=32_768)
    extension: str = Field(min_length=1, max_length=16)


class TextChunk(BaseModel):
    content: str
    chunk_index: int
    token_count: int


class DocumentChunkResponse(BaseModel):
    chunks: list[TextChunk]


class BaseParser(ABC):
    @abstractmethod
    def extract(self, path: Path) -> str:
        """Return the meaningful textual content of a document."""


class PlainTextParser(BaseParser):
    def extract(self, path: Path) -> str:
        return path.read_text(encoding="utf-8", errors="replace")


class PdfParser(BaseParser):
    def extract(self, path: Path) -> str:
        import fitz

        with fitz.open(path) as document:
            return "\n\n".join(page.get_text("text") for page in document)


class WordParser(BaseParser):
    def extract(self, path: Path) -> str:
        from docx import Document

        document = Document(path)
        blocks = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
        for table in document.tables:
            for row in table.rows:
                blocks.append(" | ".join(cell.text.strip() for cell in row.cells))
        return "\n\n".join(blocks)


class JsonParser(BaseParser):
    def extract(self, path: Path) -> str:
        parsed = json.loads(path.read_text(encoding="utf-8-sig"))
        return json.dumps(parsed, ensure_ascii=False, indent=2)


class MarkdownParser(BaseParser):
    def extract(self, path: Path) -> str:
        from lxml import html
        from markdown_it import MarkdownIt

        rendered = MarkdownIt().render(path.read_text(encoding="utf-8-sig"))
        return "\n".join(html.fromstring(rendered).itertext())


class PowerPointParser(BaseParser):
    def extract(self, path: Path) -> str:
        from pptx import Presentation

        presentation = Presentation(path)
        slides: list[str] = []
        for index, slide in enumerate(presentation.slides, start=1):
            parts = [f"Diapositiva {index}"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text.strip())
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        parts.append(" | ".join(cell.text.strip() for cell in row.cells))
            slides.append("\n".join(parts))
        return "\n\n".join(slides)


class RtfParser(BaseParser):
    def extract(self, path: Path) -> str:
        from striprtf.striprtf import rtf_to_text

        raw = path.read_text(encoding="utf-8", errors="replace")
        return rtf_to_text(raw)


class XmlParser(BaseParser):
    def extract(self, path: Path) -> str:
        from lxml import etree

        parser = etree.XMLParser(resolve_entities=False, no_network=True, recover=True)
        root = etree.parse(str(path), parser).getroot()
        return "\n".join(part.strip() for part in root.itertext() if part.strip())


class DocumentProcessor:
    """Selects a format parser, normalizes its text and creates overlapping chunks."""

    def __init__(self) -> None:
        text = PlainTextParser()
        self._parsers: dict[str, BaseParser] = {
            "pdf": PdfParser(),
            "docx": WordParser(),
            "json": JsonParser(),
            "md": MarkdownParser(),
            "txt": text,
            "pptx": PowerPointParser(),
            "rtf": RtfParser(),
            "xml": XmlParser(),
        }

    def extract_and_chunk(self, request: DocumentChunkRequest) -> DocumentChunkResponse:
        path = Path(request.path)
        if not path.is_file():
            raise ValueError("El archivo ya no existe o no es accesible")
        extension = request.extension.lower().lstrip(".")
        parser = self._parsers.get(extension)
        if parser is None:
            raise ValueError(f"El formato .{extension} no está admitido")
        text = parser.extract(path)
        normalized = re.sub(r"[ \t]+", " ", text)
        normalized = re.sub(r"\n{3,}", "\n\n", normalized).strip()
        if not normalized:
            raise ValueError("El documento no contiene texto extraíble")
        return DocumentChunkResponse(chunks=self._chunk(normalized))

    @staticmethod
    def _chunk(text: str, size: int = 1_600, overlap: int = 250) -> list[TextChunk]:
        chunks: list[TextChunk] = []
        start = 0
        while start < len(text):
            end = min(start + size, len(text))
            if end < len(text):
                boundary = max(text.rfind("\n\n", start, end), text.rfind(". ", start, end))
                if boundary > start + size // 2:
                    end = boundary + 1
            content = text[start:end].strip()
            if content:
                chunks.append(TextChunk(
                    content=content,
                    chunk_index=len(chunks),
                    token_count=max(1, len(content) // 4),
                ))
            if end >= len(text):
                break
            start = max(start + 1, end - overlap)
        return chunks
